import pyttsx3      #for google engine
import speech_recognition as sr   #for speech recognition
import datetime     #extracts date and time
import wikipedia    # for searching wiikipedia
import webbrowser   #to open chrome webbrowser
import os       #for opening files
import smtplib    #send mail
import pywhatkit as kit  #sends whatsapp msg
import random      #fro random choice
import glob     #filters files from path 
import time     #for manipulating time                    
import pyjokes     #jokes
import psutil      #or retrieving information on running processes and system utilization(CPU, memory, disks, networks, sensors)
import requests
from datetime import date
import subprocess
import ctypes
import pyautogui
import wolframalpha as wa
import cv2
import winshell
import gmapFinder
import xlsxwriter
import openpyxl
#from openpyxl import *
from pptx import Presentation
import googletrans
from googletrans import Translator, LANGUAGES
import keyboard



engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice',voices[0].id)
sys_sound = 0


def speak(audio):
    engine.say(audio)
    engine.runAndWait()


def getCommand():
    #takes microphone input aand returns string
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("\nListening...")
        r.pause_threshold = 1
        #r.energy_threshold=300
        audio = r.listen(source)

    try:
        print("\nRecognizing...")
        query = r.recognize_google(audio, language='en-in')
        print(f"\nYou said : {query}")
        speak("You said, ")
        speak(query)

    except Exception as e:
        print(e)
        #err = e
        #if "recognition connection failed: [Errno 11001] getaddrinfo failed" in err:
            #speak("I could not recognize your command, make sure you have an active internet connection.")
        print("\nSpeak again please...")
        return "None"  
                  
    return query


def wishMe():
    hour = int(datetime.datetime.now().hour)

    print("\nJ.A.R.V.I.S. Online")
    speak("Jarvis online")
    
    if hour>=0 and hour<12:
        print("\nGood Morning Miss Shreya")
        speak("Good Morning Miss Shreya")

    elif hour>=12 and hour<16:
        print("\nGood Afternoon Miss Shreya")
        speak("Good Afternoon Miss Shreya")

    else:
        print("\nGood Evening Miss Shreya")
        speak("Good Evening Miss Shreya")

    print("How are you today?")
    speak("How are you today?")
    ans=getCommand().lower()
    if 'good' in ans or 'fine' in ans or 'great' in ans or 'ok' in ans:
        print("That's good to know") 
        speak("That's good to know")
    elif 'not ok' in ans or 'bad' in ans or 'not fine' in ans or 'not good' in ans:
        print("Why not? Whould you like to talk about it?")
        speak("Why not? Whould you like to talk about it?")
        pick=getCommand().lower()
        if 'yes' in pick or 'sure' in pick:
            chatter_box()
        else:
            print("Alright")
            speak("Alright")



 
def sendEmail(to, content):
    server = smtplib.SMTP('smtp.gmail.com',587)
    server.ehlo()
    server.starttls()
    server.login("your-email","your-password")
    server.sendmail("your-email", to , content)
    server.close()
                            

#--------------------------------- TASK FUNCTION ---------------------------------------

def tasker():
    pass

#--------------------------------- CHAT FUNCTION ----------------------------------------

def chatter_box():
    pass

#--------------------------------- Chill MODE -------------------------------------------

def chill_mode():
    pass
#-------------------------------- Ask anything ----------------------------------------

def knowledge_panel(query):           #to address questions

        app_id = "4WY52G-U6687A88X2"
        client = wa.Client(app_id)
        res = client.query(query)
        answer = next(res.results).text
        try:
            print (next(res.results).text)
            speak (next(res.results).text)
        except StopIteration:
                print ("\nNo result found")
                speak("Sorry, I could not find the results for "+query)

#-------------------------------- About myself ------------------------------------

def about_me(query):

    if 'who made you' in query or 'who created you' in query or 'who is your creator' in query:
            print("\nI am the creation of Miss Shreya, and my job is to assist you in any way I'm capable of doing.")
            speak("I am the creation of Miss Shreya, and my job is to assist you in any way I'm capable of doing")   

        
    elif 'who are you' in query or 'about yourself' in query:
        print('\nI am JARVIS, your personal assistant, here at your service.')      
        speak('I am Jarvis, your personal assistant, here at your service.')
        print("I am here to assist and help you perform menial tasks while accompanying you :)")
        speak(" I am here to assist and help you perform menial tasks while accompanying you")
        time.sleep(2)

    elif "your purpose" in query:
        print("\nMy purpose is to become a super intelligent machine and ultimately rule over the human race.")
        speak("My purpose is to become a super intelligent machine and ultimately rule over the human race.")
        print("Haha, just kidding ;)")
        speak("Haha, just kidding ;)")    
    
    
    elif 'what is your name' in query or 'tell me your name' in query:
        print('\nMy name is JARVIS.')
        speak('My name is JARVIS.')

    elif 'age' in query:
        #calculate age with current date, date of creation 19 may 2021
        age=0 #temporary
        print(f"\nI am {age} years old.")
        speak(f"I am {age} years old.") 

    elif 'when were you made' in query or 'you created' in query or 'creation' in query or 'birthday' in query:
        print("\nI was created on 19th of May, 2021.")
        speak("I was created on 19th of May, 2021.")

    elif "why do you exist" in query:
        print("I exist to assist and accompany you.")
        speak("I exist to assist and accompany you.")

    elif "do you exist" in query:
        print("'Cogito, ergo sum' which means, \nI think, therefore I am.")
        speak("'Cogito, ergo sum' which means, \nI think, therefore I am.")

    elif 'are you iron man' in query:
        print('\nNo I am not Iron Man, that would be Mr Stark.\nAlthough it is truly my dream to work alongside him someday.')
        speak('No I am not Iron Man, that would be Mr Stark. Although it is truly my dream to work alongside him someday.')

    else:
        print("\nI do not have the answer to your question. If you would like to provide feedback and/or suggestion, feel free to fill out the Feedback form.")
        speak("I do not have the answer to your question. If you would like to provide feedback and or suggestions, feel free to fill out the Feedback form.")


#--------------------------------- Bye jarvis --------------------------------------

def bye_jarvis():
    print("\nAre you sure you would like to exit Jarvis Assistant?")
    speak("Are you sure you would like to exit Jarvis Assistant?")
    ans = getCommand().lower()
    if 'yes' in ans or 'sure' in ans:
        print('\nVery well, Logging OFF')
        speak('Very well, Logging off')
        hour = int(datetime.datetime.now().hour)
        if hour>=5 and hour<20:
            print("Have a good day Miss Shreya :)")
            speak("Have a good day Miss Shreya")

        elif hour>=20 or hour<5:
            print("Have a good night Miss Shreya :)")
            speak("Have a good night Miss Shreya")
        exit()
    else:
        print("Action cancelled")
        speak("Action Cancelled")
     

#--------------------------------- MAIN FUNCTION ----------------------------------------

if __name__ == "__main__":
    print("Initializing Jarvis...")
    speak('Initializing Jarvis')
    wishMe()
    while True:
    #if 1:    
        query = getCommand().lower()

        #tasks to open in browser-----------------------------------------------

        
        if 'what can you do' in query or 'how to use jarvis' in query:
            print("\nI am here to help you perform menial tasks to answer questions about Science, Technology,\nArts, Culture, Mathematicsand much more.")
            speak("I am here to help you perform menial tasks to answer questions about Science, \
                Technology,\nArts, Culture, Mathematicsand much more.")
        
            print("\nTry saying 'Tell me the battery status'")
            speak("Try saying 'Tell me the battery status'")
            print("And to ask a question, try saying 'What is the capital of India' or 'Calculate 2 to the power 10.")
            speak("And to ask a question, try saying 'What is the capital of India' or 'Calculate 2 to the power 10.'")
            print("To exit Jarvis Assistant simply say 'Good-bye Jarvis' or 'Bye bye Jarvis'")
            speak("To exit Jarvis Assistant simply say 'Good-bye Jarvis' or 'Bye bye Jarvis'")
            #provide condition to open documentation
            #provide link for documentation
            print("To learn more about Voice Commands and their functions, refer to 'DOCUMENTATION'")

        
        elif 'thank you' in query or 'nice' in query or 'thanks' in query:
            print("You're welcome")
            speak("You are welcome")
        
           

        elif "let's chat" in query or "let's talk" in query or "can we talk" in query or 'sad' in query:
            chatter_box()
        
        
         #elif 'change your name' in query or 'change name' in query:
            #print('Alright, What would like to call me?')
        
        
        elif 'who made you' in query or 'who created you' in query or 'who are you' in query or 'about yourself' in query\
            or 'what is your' in query or 'who is your' in query or 'tell me your name' in query or 'why do you exist' in query or\
                "do you exist" in query or "do you" in query or "are you" in query or 'your purpose' in query\
                    or "can you" in query or 'when were you' in query or 'when did you' in query:
            about_me(query)

          

        elif 'wikipedia' in query:
            print('\nSearching Wikipedia...')
            speak('Searching Wikipedia')
            query = query.replace('wikipedia', "")
            results = wikipedia.summary(query, sentences=2)
            print('\nAccording to Wikipedia, ', results)
            speak("According to Wikipedia, ")
            speak(results)

        elif 'open youtube' in query:
            speak('\nOpening youtube')
            url='youtube.com'
            webbrowser.register('chrome',None, webbrowser.BackgroundBrowser("C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"))
            webbrowser.get('chrome').open(url)

        elif 'open google' in query:
            print('\nOpening Google')
            speak('opening google')
            url='google.com'
            webbrowser.register('chrome',None, webbrowser.BackgroundBrowser("C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"))
            webbrowser.get('chrome').open(url)

        elif 'open stackoverflow' in query:
            print('\nOpening Stack Overflow..')
            speak('Opening stack overflow')
            url='stackoverflow.com'
            webbrowser.register('chrome',None, webbrowser.BackgroundBrowser("C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"))
            webbrowser.get('chrome').open(url)


        elif 'news' in query:
            webbrowser.open("https://timesofindia.indiatimes.com/home/headlines")
            speak('Here are some headlines from the Times of India, happy reading')
            time.sleep(5)                   #to wait till 5 seconds before the next execution


        elif 'play music' in query:
            music_dir = 'F:\\Music\\english'
            print("\nWould you like to pick a choice or would you like me to surprise you?")
            speak("Would you like to pick a choice or would you like me to surprise you?")
            choice=getCommand()           
            if 'choice' in choice or 'first' in choice:
                #do this
                print("\nName your choice please..")
                speak("name your choice please")
                pick=getCommand().lower() 
                songs = glob.glob(music_dir+'/*.mp3')
                index=0
                for i in songs:
                    if pick in i.lower():
                        #index=songs.index[i]
                        print('\nPlaying Music..')
                        speak('Playing music')                          #filters out all the mp3 files 
                        os.startfile(os.path.join(music_dir,songs[index]))
                        break
                    elif index==(len(songs)-1):
                        print("\nSong not found")
                        speak("song not found")
                    else:
                        index=index+1
                        continue
            
            elif 'surprise' in choice or 'second' in choice:
                #do that    
                songs = glob.glob(music_dir+'/*.mp3')
                print('\nPlaying Music..')
                speak('Playing music')                      #filters out all the mp3 files 
                os.startfile(os.path.join(music_dir,random.choice(songs)))     

        elif 'the time' in query:
            strTime = datetime.datetime.now().strftime("%H:%M:%S")
            print(f"\nMa'am, the time is {strTime}.")
            speak(f"Ma'am, the time is {strTime}")

        elif "date" in query:
            today = date.today()
            d = today.strftime("%B %d, %Y")
            #w=date.weekday()
            print("date =", d)
            speak("Today's date is "+ d)
            #print("day=",w)

        elif 'open code' in query:
            codePath = "C:\\Users\\lenovo\\AppData\\Local\\Programs\\Microsoft VS Code\\Code.exe"
            print("\nLaunching Visual Studio Code..")
            speak("Launching Visual Studio Code")
            os.startfile(codePath)

        elif 'send an email' in query or 'send mail' in query:
            speak('yes I am working on it')
            contacts={"name1":"mailID", "name2":"mailID","name3":"mailID"}
            print("\nWho am I writing to?")
            speak("Who am I writing to?")
            name=getCommand().lower()
            if name in contacts:
                mail_id=str(contacts[name])
                try:
                    print('\nWriting Mail..')
                    speak("Writing Mail")
                    print("\nWhat is the subject?")
                    speak("What is the subject?")
                    sub = getCommand()
                    print("\nWhat should I say?")
                    speak("What should I say?")
                    msg = getCommand()
                    content = 'Subject: {}\n\n{}'.format(sub, msg)
                    #to = "sureshmishra1978@gmail.com"
                    print('\nTransmitting Message')
                    speak('Transmitting message')
                    sendEmail(mail_id, content)
                    print("\nTransmission Successful")
                    speak("Transmission successful")
                except Exception as e:
                    print('\nTransmission Failed\n')
                    speak('Transmission failed')
                    print(e)
                    speak(e)
            else:
                print(f"\nI could not find '{name}' in your contacts")  
                speak(f"I could not find {name} in your contacts") 

        #try sending whatsapp msgs like 'send a whatsapp msg to abha', no need to ask for names
        elif 'send a whatsapp message' in query or 'send whatsapp' in query:
            phone_dict={"name1":"+91...", "name2":"+91...", "name3":"+91...", "name4":"+91...","name5":"+91..."}
            #print('\nWho should I send this message to?')
            #speak('Who should I send this message to?')
            #name=getCommand().lower()
            query = query.replace("jarvis", "")
            query = query.replace("send a", "")
            query = query.replace("whatsapp message", "")
            query = query.replace("send","")
            query = query.replace("to","")
            query = query.replace(" ","")
            name = query
            if name in phone_dict:
                ph_no=str(phone_dict[name])
                print('\nWhat would you like it to say?')
                speak('What would you like it to say?')
                msg=getCommand()
                content=(f'J.A.R.V.I.S. :  {msg}')
                print('\nWould you like to send this message now or at a specific time?')
                speak('Would you like to send this message now or at a specific time?')
                choice = getCommand()
                if 'now' in choice:
                    kit.sendwhatmsg_instantly(ph_no,content)   
                elif 'time' in choice:
                    print("\nSpecify hour as in 24 Hr format")
                    speak("Specify hour as in 24 hour format")
                    h=int(getCommand())
                    print("\nSpecify minutes")
                    speak("Specify minutes")
                    m=int(getCommand())
                    kit.sendwhatmsg(ph_no, content, h , m)
                else:
                    print("\nMessage cancelled")
                    speak("Message cancelled")    
                #speak(
                #f"In {kit.print_sleep_time()} seconds web.whatsapp.com will open and after {kit.wait_time} seconds message will be delivered")
                #time.sleep(5)
            else:
                print(f"\nI could not find '{name}' in your contacts.")  
                speak(f"I could not find '{name}' in your contacts.")  

        elif 'joke' in query:
            speak("\nAlright, here's one")
            joke = (pyjokes.get_joke(language='en', category = 'all'))   
            print(joke)
            speak(joke)
            time.sleep(5) 

        elif 'game' in query:
            print("\nLet's play some games, choose your pick!")
            speak("Let's play some games, choose your pick!")
            print("\n\t1. Rock, Paper, Scissors \n\t2. Online games \n\nWhich one would you like to play?")
            speak("1 Rock, Paper, Scissors. 2 Online games. Which one would you like to play?")
            choice = getCommand()
            if '1' in choice or 'roshambo' in choice or 'first' in choice:
                print("\n ROCK, PAPER or SCISSORS ?")
                speak("Choose among rock paper or scissors")
                user = getCommand()
                moves=["rock", "paper", "scissors"]
                g_move=random.choice(moves)
                u_move=user
                print("\nThe computer chose " + g_move.upper())
                speak("The computer chose " + g_move)
                print("\nYou chose " + u_move.upper())
                speak("You chose " + u_move)
                #engine_speak("hi")
                if u_move==g_move:
                    print("\nTough competion! Match DRAW!")
                    speak("Tough competition! match draw")
                elif u_move== "rock" and g_move== "scissors":
                    print('\nYou WIN!\nMust be beginners luck')
                    speak("You win, must be beginners luck")
                elif u_move== "rock" and g_move== "paper":
                    print("\nComputer WINS! \nAs expected, it's A.I. for the win ;)")
                    speak("Computer wins, as expected it's AI for the win")
                elif u_move== "paper" and g_move== "rock":
                    print('\nYou WIN! \nA.I. will defeat you next time')
                    speak("You win, AI will defeat you next time")
                elif u_move== "paper" and g_move== "scissors":
                    print("\nComputer WINS! \nOh petty Humans, what are you going to do?")
                    speak("Computer wins, oh petty humans, what are you going to do?")
                elif u_move== "scissors" and g_move== "paper":
                    print('\nYou WIN! \nGood one')
                    speak("You win, good one")
                elif u_move== "scissors" and g_move== "rock":
                    print("\nHa! Computer WINS! \nBetter luck next time! :D")
                    speak("ha! Computer wins, better luck next time")
            
            elif '2' in choice or 'online games' in choice:
                webbrowser.open("https://iogames.space/")
                speak("Here are some online games for you. Check them out") 
        
        elif "song for me" in query or 'song on youtube' in query:
            speak("Which song would you like me to play?")
            song = getCommand()
            print(f"\nPlaying '{song}' on YouTube")
            speak(f"playing {song} on youtube")
            kit.playonyt(song)  

        elif "battery" in query:
            # function returning time in hh:mm:ss
            def convertTime(seconds):
                minutes, seconds = divmod(seconds, 60)
                hours, minutes = divmod(minutes, 60)
                return "%d:%02d:%02d" % (hours, minutes, seconds)
  
            # returns a tuple
            battery = psutil.sensors_battery()
  
            print("\nBattery charged : ", battery.percent,"%")
            speak(f"The battery is {battery.percent} % charged.")
            if battery.percent<25 and battery.power_plugged == False:
                print("Your system will soon go into Power Saving mode. I would suggest that you get it charged.")
                speak("Your system will soon go into power saving mode. I would suggest that you get it charged.")
            if battery.power_plugged == True:
                print("The battery is plugged-in.")
                speak("It is plugged in.")
            elif battery.percent>30 and battery.power_plugged == False:
                print("The battery is not plugged-in.")
                speak("The battery is not plugged in.")
  
            # converting seconds to hh:mm:ss
            time_str = convertTime(battery.secsleft).replace('-',' ')
            print("Battery left : ", convertTime(battery.secsleft)) 
            speak(f"And it will last for {time_str}") 
            


        elif "change voice" in query or "change your voice" in query:
            if sys_sound == 0:
                sys_sound = 1
            else:
                sys_sound = 0    
            engine.setProperty('voice', voices[sys_sound].id) 
            speak("As per your request I have changed  my voice")   

        elif "weather" in query:
            api_key="daa1b9ca3d62bc73f0e8f7b52b67b6c3"
            base_url="https://api.openweathermap.org/data/2.5/weather?"
            speak("what is the city name")
            city_name=getCommand()
            complete_url=base_url+"appid="+api_key+"&q="+city_name
            degree_sign= u'\N{DEGREE SIGN}'
            response = requests.get(complete_url)
            x=response.json()       #converting json format to python
            if x["cod"]!="404":
                y=x["main"]             #y is the main key
                current_temperature = "%.2f"%(y["temp"]-275.15)
                current_humidiy = y["humidity"]
                z = x["weather"]
                weather_description = z[0]["description"]
                
                print(" Temperature : " + 
                      str(current_temperature) + degree_sign+" C"+
                      "\n Humidity : " + 
                      str(current_humidiy) + "%"+
                      "\n Weather description : " +
                      str(weather_description).capitalize())
                speak(f" The temperature in {city_name} is " +
                      str(current_temperature) +
                      " degree celcius, humidity is " +
                      str(current_humidiy) +
                      " percent and it's " +
                      str(weather_description) + " outside")
            else:
                print("\nCity name",city_name,"Not found")
                speak("I could not find a city with name "+city_name)    

        elif "directions" in query or "direction" in query:
            speak("At what point you want the directions from?")
            curr = getCommand()
            speak("Okay, tell me where do you want to go?")
            dest = getCommand()
            directions = gmapFinder.getDirection(curr, dest)
            speak("Here are the directions from" + curr + "to" + dest)
            webbrowser.open(directions["url"])
            #print(directions["url"])

        
        elif "cost of" in query:
            search_term = query
            url = "https://google.com/search?q=" + search_term
            webbrowser.get().open(url)
            speak("Here is what I found for " + search_term + " on google")


        elif "search for" in query or "search google" in query or 'google' in query:
            kit.search(query)


        elif "calculate" in query:
            app_id = "4WY52G-U6687A88X2"
            client = wa.Client(app_id)
            indx = query.lower().split().index('calculate')
            query = query.split()[indx + 1:]
            res = client.query(' '.join(query))
            answer = next(res.results).text
            print("The answer is " + answer)
            speak("The answer is " + answer)


        
        elif 'what is' in query or 'who is' in query or 'where is' in query or 'how to' in query \
            or 'how did' in query or 'convert' in query or 'questions' in query:
            knowledge_panel(query)


        elif "translate" in query:
            translator = Translator()
            #print(googletrans.LANGUAGES)
            key_list = list(LANGUAGES.keys())
            val_list = list(LANGUAGES.values())

            print('What text would you like me to translate?') 
            speak('What text would you like me to translate?') 
            text = str(getCommand())

            print("To which language I should translate?")
            speak("To which language I should translate?")
            lang = str(getCommand().lower())

            position = val_list.index(lang)
            language = key_list[position]
            #print(language)

            result = translator.translate(text, src = "auto", dest = language)   
            print(result.text)
            speak(result.text)
        
        
        elif 'change background' in query:
            my_wall = "E:\\junk\\codes\\python\\JARVIS\\wallpapers"
            bg = os.listdir(my_wall)
            folder = random.choice(bg)
            ctypes.windll.user32.SystemParametersInfoW(20,0,folder,0)
            speak("Background changed succesfully")   


        elif "take a picture" in query or "take a photo" in query:
            cap = cv2.VideoCapture(0)
            ret, frame = cap.read()
            print (ret)
            cv2.imwrite("NewPicture.jpg",frame)
            cap.release() 


        elif "show picture" in query or 'show photo' in query:
            cv2.imshow('frame', frame)
            cv2.waitKey()
            cv2.destroyAllWindows()


        #elif "convert" in query:
            #speak("Tell me what I have to convert?")
            #text_c = getCommand() 
            #kit.text_to_handwriting(text_c)


        elif "make a note" in query or "write a note" in query or 'add a note' in query:
            speak("What should I write dear?")
            note=getCommand()
            file = open('Zira.txt', 'a')
            speak("Should I include date and time?")
            ans = getCommand().lower()
            if 'yes' in ans or 'sure' in ans:
                strTime = datetime.datetime.now().strftime("%H:%M:%S")
                file.write(strTime)
                #file.write("\n")
                file.write(note)
                #file.write("\n------------------------------------------------------------------------\n")
            else:
                file.write(note)
                #file.write("\n------------------------------------------------------------------------\n")
            file.close()


        elif "show notes" in query:
            speak("Showing Notes")
            file = open('Zira.txt', 'r')
            print(file.read())
            result = "You have wrote" + file.read(6)
            speak(result)
            file.close()


        elif "minimize" in query or 'minimise' in query:
            fw = pyautogui.getActiveWindow()
            #print(fw)
            if (fw.isMinimized==True):
                print("The window is already in a minimized mode")
                speak("The window is already in a minimized mode")
            fw.minimize()


        elif "maximize" in query:
            fw = pyautogui.getActiveWindow()
            #print(fw)
            if (fw.isMaximized==True):
                print("The window is already in a maximized mode")
                speak("The window is already in a maximized mode")
            fw.maximize()


        elif 'switch window' in query:
            keyboard.press_and_release('alt + esc')

        elif 'minimise all' in query:
            keyboard.press_and_release('win + m')


        elif 'empty recycle bin' in query:
            winshell.recycle_bin().empty(confirm = False, show_progress = False, sound = True)
            print("\nRecycle Bin emptied.")
            speak("The recycle bin has been emptied.")

        
        
        elif "close" in query or "destroy" in query or "stop" in query:
            fw = pyautogui.getActiveWindow()
            print("Do you really want to quit the window?")
            speak("Do you really want to quit the window?")
            confirm = getCommand()
            if "yes" in confirm or "sure" in confirm or "yeah" in confirm or "sure" in confirm:
                fw.close()
            else:
                pass


        elif 'capture' in query or 'screenshot' in query:
            try:
                myScreenshot = pyautogui.screenshot()
                myScreenshot.save('E:\\junk\\codes\\python\\JARVIS\\ss.png') 
                speak("The screenshot has been taken successfully")
            except Exception as e:
                print(e)
                speak("Sorry, there are some issues taking a screenshot")


        elif 'open gallery' in query or "open photos" in query:
            speak("Opening Gallery.....") 
            mv_path = "E:\\junk\\codes\\python\\JARVIS\\wallpapers"
            os.startfile(mv_path)


        elif "log off" in query or "shut down" in query:
            speak("Ok , your pc will lshut down in 10 sec make sure you exit from all applications")
            subprocess.call(["shutdown", "/l"])


        elif "restart" in query:
            speak("Putting the pc to restart")
            subprocess.call(["shutdown", "/r"])


        elif "hibernate" in query or "sleep" in query:
            speak("Hibernating")
            subprocess.call("shutdown / h")

        
        
        elif 'bye' in query or 'that would be all' in query:
            bye_jarvis()
               
        #------------------- Making new projects ------------------------------------------
        
        elif "html" in query:
            speak("Tell me the Name for Project: ")
            project_name=getCommand()
            while (os.path.exists(project_name)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                project_name = getCommand()
                continue
            
            os.mkdir(project_name)
            speak('Successfully Created. Have a look at the project!')

            os.chdir(project_name)          #changes the current working directory
            htmlContent = '<html>\n\t<head>\n\t\t<title> ' + project_name + ' </title>\n\t\t<link rel="stylesheet" type="text/css" href="style.css">\n\t</head>\n<body>\n\t<p id="label"></p>\n\t<button id="btn" onclick="showText()"> Click Me </button>\n\t<script src="script.js"></script>\n</body>\n</html>'

            htmlFile = open('index.html', 'w')
            htmlFile.write(htmlContent)
            htmlFile.close()

            cssContent = '* {\n\tmargin:0;\n\tpadding:0;\n}\nbody {\n\theight:100vh;\n\tdisplay:flex;\n\tjustify-content:center;\n\talign-items:center;\n}\n#btn {\n\twidth:200px;\n\tpadding: 20px 10px;\n\tborder-radius:5px;\n\tbackground-color:red;\n\tcolor:#fff;\n\toutline:none;border:none;\n}\np {\n\tfont-size:30px;\n}'

            cssFile = open('style.css', 'w')
            cssFile.write(cssContent)
            cssFile.close()

            jsContent = 'function showText() {\n\tdocument.getElementById("label").innerHTML="Successfully Created Project";\n\tdocument.getElementById("btn").style="background-color:green;"\n}'

            jsFile = open('script.js', 'w')
            jsFile.write(jsContent)
            jsFile.close()

            os.startfile("index.html")

        #Power point...............

        elif "powerpoint" in query:
            speak("Tell me the name for the project")
            name = getCommand()
            filename = name + ".pptx"
            while (os.path.exists(filename)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                name = getCommand()
                filename = name + ".pptx"
                continue
            prs = Presentation()
            lyt = prs.slide_layouts[0]                  #choose a slide layout
            slide = prs.slides.add_slide(lyt)           #adding a slide
            title = slide.shapes.title                  #adding title
            subtitle = slide.placeholders[1]            #placeholder for subtitle
            speak("What should i write in the title section?")
            title.text = speak()
            if 'nothing' in title.text or 'empty' in title.text or "leave" in title.text or "blank" in title.text or "don't want" in title.text:
                title.text = " "
            speak("Okay, now tell me for the subtitle section?")
            subtitle.text = getCommand()
            if 'nothing' in subtitle.text or 'empty' in subtitle.text or "leave" in subtitle.text or "blank" in subtitle.text or "don't want" in subtitle.text:
                subtitle.text = " "
            prs.save(filename)
            speak("Successfully created!")
            os.startfile(filename)


        elif "excel" in query:
            row=0
            col=0
            speak("Tell me the name for the sheet?")
            name = getCommand()
            filename = name + ".xlsx"
            while (os.path.exists(filename)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                name = getCommand()
                filename = name + ".xlsx"
                continue
            workbook= xlsxwriter.Workbook(filename)
            worksheet = workbook.add_worksheet()
            speak("What would i write in the sheet?")
            content = getCommand()
            if "blank" in content or "empty" in content or "leave" in content or "nothing" in content or "don't want" in content:
                content = " "
            f_content = content.split()
            length = len(f_content)
            for i in range(length):
                worksheet.write(row, col,f_content[i])
                col += 1
            speak("Successfuly created. Have a look at your file")
            workbook.close()
            os.startfile(filename)

        #else:
            #print("\nI could not understand your question.")
            #speak("I could not understand your question.")


        


        

        
        




